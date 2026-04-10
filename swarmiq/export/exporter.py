"""
Export_Module for SwarmIQ v2.
Renders research reports in Markdown, PDF, PPT, and JSON formats.
"""
from __future__ import annotations

import io
import json
import dataclasses
from typing import Union

from swarmiq.core.models import Reference, Figure
from swarmiq.export.citations import format_apa, format_mla


class ExportError(Exception):
    """Raised when a PDF or PPT rendering operation fails."""


class ExportModule:
    """Renders the final report in multiple formats with citation support."""

    SUPPORTED_FORMATS = {"markdown", "pdf", "ppt", "json"}

    def export(
        self,
        report_markdown: str,
        references: list[Reference],
        figures: list[Figure],
        format: str,
        citation_style: str = "apa",
    ) -> Union[bytes, str]:
        """Export the report in the requested format.

        Args:
            report_markdown: The Markdown text of the report.
            references: List of Reference objects for the citation list.
            figures: List of Figure objects (Plotly JSON or matplotlib PNG bytes).
            format: One of "markdown", "pdf", "ppt", "json".
            citation_style: "apa" or "mla" (ignored for JSON).

        Returns:
            bytes for PDF/PPT, str for Markdown/JSON.

        Raises:
            ValueError: If an unsupported format is requested.
            ExportError: If PDF or PPT rendering fails.
        """
        fmt = format.lower()
        if fmt not in self.SUPPORTED_FORMATS:
            raise ValueError(
                f"Unsupported format '{format}'. Choose from: {', '.join(sorted(self.SUPPORTED_FORMATS))}"
            )

        if fmt == "markdown":
            return self._export_markdown(report_markdown, references, citation_style)
        if fmt == "pdf":
            return self._export_pdf(report_markdown, references, figures, citation_style)
        if fmt == "ppt":
            return self._export_ppt(report_markdown, references, figures, citation_style)
        # json
        return self._export_json(report_markdown, references, figures)

    # ── Format implementations ────────────────────────────────────────────────

    def _export_markdown(
        self,
        report_markdown: str,
        references: list[Reference],
        citation_style: str,
    ) -> str:
        ref_text = self._format_references(references, citation_style)
        if ref_text:
            return f"{report_markdown}\n\n## References\n\n{ref_text}"
        return report_markdown

    def _export_pdf(
        self,
        report_markdown: str,
        references: list[Reference],
        figures: list[Figure],
        citation_style: str,
    ) -> bytes:
        try:
            from fpdf import FPDF  # fpdf2 exposes fpdf.FPDF
        except ImportError as exc:
            raise ExportError("fpdf2 is not installed. Install it with: pip install fpdf2") from exc

        try:
            pdf = FPDF()
            pdf.set_margins(left=20, top=20, right=20)
            pdf.set_auto_page_break(auto=True, margin=20)
            pdf.add_page()
            pdf.set_font("Helvetica", size=11)

            # Effective content width
            w = pdf.w - pdf.l_margin - pdf.r_margin

            # Write report text line by line
            for line in report_markdown.splitlines():
                # Render headings slightly larger
                if line.startswith("## "):
                    pdf.set_font("Helvetica", style="B", size=13)
                    pdf.multi_cell(w, 8, line[3:])
                    pdf.set_font("Helvetica", size=11)
                elif line.startswith("# "):
                    pdf.set_font("Helvetica", style="B", size=15)
                    pdf.multi_cell(w, 10, line[2:])
                    pdf.set_font("Helvetica", size=11)
                else:
                    pdf.multi_cell(w, 6, line if line else " ")

            # Embed matplotlib PNG figures
            for fig in figures:
                if fig.figure_type == "matplotlib" and isinstance(fig.data, bytes) and fig.data:
                    pdf.add_page()
                    pdf.set_font("Helvetica", style="B", size=11)
                    pdf.cell(w, 8, f"Figure: {fig.figure_id}", new_x="LMARGIN", new_y="NEXT")
                    img_buf = io.BytesIO(fig.data)
                    pdf.image(img_buf, w=min(w, 170))

            # Append reference list
            ref_text = self._format_references(references, citation_style)
            if ref_text:
                pdf.add_page()
                pdf.set_font("Helvetica", style="B", size=13)
                pdf.multi_cell(w, 8, "References")
                pdf.set_font("Helvetica", size=10)
                for line in ref_text.splitlines():
                    pdf.multi_cell(w, 6, line if line else " ")

            return bytes(pdf.output())
        except ExportError:
            raise
        except Exception as exc:
            raise ExportError(f"PDF rendering failed: {exc}") from exc

    def _export_ppt(
        self,
        report_markdown: str,
        references: list[Reference],
        figures: list[Figure],
        citation_style: str,
    ) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ExportError(
                "python-pptx is not installed. Install it with: pip install python-pptx"
            ) from exc

        try:
            prs = Presentation()
            slide_layout_title = prs.slide_layouts[0]  # Title Slide
            slide_layout_content = prs.slide_layouts[1]  # Title and Content

            # Split report into sections by ## headings
            sections = self._split_sections(report_markdown)

            # Title slide — use first heading or generic title
            title_slide = prs.slides.add_slide(slide_layout_title)
            title_slide.shapes.title.text = sections[0]["heading"] if sections else "Research Report"
            if title_slide.placeholders[1]:
                title_slide.placeholders[1].text = "Generated by SwarmIQ v2"

            # One slide per section
            for section in sections:
                slide = prs.slides.add_slide(slide_layout_content)
                slide.shapes.title.text = section["heading"]
                tf = slide.placeholders[1].text_frame
                tf.word_wrap = True
                # Truncate body to avoid overflow
                body_text = section["body"][:1500]
                tf.text = body_text

            # Embed matplotlib figures
            for fig in figures:
                if fig.figure_type == "matplotlib" and isinstance(fig.data, bytes) and fig.data:
                    fig_slide = prs.slides.add_slide(slide_layout_content)
                    fig_slide.shapes.title.text = f"Figure: {fig.figure_id}"
                    img_buf = io.BytesIO(fig.data)
                    prs.slides[-1].shapes.add_picture(
                        img_buf, Inches(1), Inches(1.5), width=Inches(8)
                    )

            # References slide
            ref_text = self._format_references(references, citation_style)
            if ref_text:
                ref_slide = prs.slides.add_slide(slide_layout_content)
                ref_slide.shapes.title.text = "References"
                tf = ref_slide.placeholders[1].text_frame
                tf.word_wrap = True
                tf.text = ref_text[:2000]

            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()
        except ExportError:
            raise
        except Exception as exc:
            raise ExportError(f"PPT rendering failed: {exc}") from exc

    def _export_json(
        self,
        report_markdown: str,
        references: list[Reference],
        figures: list[Figure],
    ) -> str:
        def _ref_to_dict(ref: Reference) -> dict:
            return {
                "ref_id": ref.ref_id,
                "url": ref.url,
                "title": ref.title,
                "authors": ref.authors,
                "year": ref.year,
            }

        def _fig_to_dict(fig: Figure) -> dict:
            data = fig.data
            # bytes → base64 string for JSON serialisation
            if isinstance(data, bytes):
                import base64
                data = base64.b64encode(data).decode("utf-8")
            return {
                "figure_id": fig.figure_id,
                "figure_type": fig.figure_type,
                "data": data,
            }

        payload = {
            "report_markdown": report_markdown,
            "references": [_ref_to_dict(r) for r in references],
            "figures": [_fig_to_dict(f) for f in figures],
        }
        return json.dumps(payload, indent=2)

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _format_references(self, references: list[Reference], citation_style: str) -> str:
        if not references:
            return ""
        style = citation_style.lower()
        if style == "mla":
            return format_mla(references)
        return format_apa(references)

    def _split_sections(self, markdown: str) -> list[dict]:
        """Split Markdown into sections by ## headings."""
        sections: list[dict] = []
        current_heading = "Introduction"
        current_body_lines: list[str] = []

        for line in markdown.splitlines():
            if line.startswith("## "):
                if current_body_lines or sections:
                    sections.append(
                        {"heading": current_heading, "body": "\n".join(current_body_lines).strip()}
                    )
                current_heading = line[3:].strip()
                current_body_lines = []
            elif line.startswith("# "):
                # Top-level heading becomes the first section title
                if current_body_lines or sections:
                    sections.append(
                        {"heading": current_heading, "body": "\n".join(current_body_lines).strip()}
                    )
                current_heading = line[2:].strip()
                current_body_lines = []
            else:
                current_body_lines.append(line)

        # Flush last section
        sections.append(
            {"heading": current_heading, "body": "\n".join(current_body_lines).strip()}
        )
        return sections
